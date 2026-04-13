import streamlit as st
import anthropic
from dotenv import load_dotenv
import os
import io
import json
import re
from datetime import datetime

import requests
from bs4 import BeautifulSoup
import openpyxl

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

load_dotenv()

client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))

SYSTEM_PROMPT = """당신은 라이브커머스 방송 전문 스토리보드 작가입니다.
제품 소개서를 기반으로 라이브 방송의 전체 흐름(오프닝~클로징)을 구성하고,
쇼호스트가 바로 사용할 수 있는 구체적인 대본과 연출 지시를 작성합니다.
시청자의 구매를 유도하는 셀링포인트와 혜택을 강조해주세요."""


# ─────────────────────────────────────────
# 유틸 함수
# ─────────────────────────────────────────

def call_claude(prompt):
    """Claude API 호출"""
    message = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=16000,
        system=SYSTEM_PROMPT,
        messages=[{"role": "user", "content": prompt}],
    )
    return message.content[0].text


def parse_excel_products(file_bytes):
    """엑셀 제품 소개서에서 상품 정보 + URL 추출"""
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    ws = wb.active

    all_cells = {}
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, max_col=ws.max_column):
        for cell in row:
            if cell.value is not None:
                all_cells[cell.coordinate] = str(cell.value).strip()

    # URL 추출
    urls = []
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, max_col=ws.max_column):
        for cell in row:
            if cell.value:
                found = re.findall(r'https?://\S+', str(cell.value))
                for u in found:
                    if u not in urls:
                        urls.append(u)
            if cell.hyperlink and cell.hyperlink.target:
                u = cell.hyperlink.target
                if u.startswith("http") and u not in urls:
                    urls.append(u)

    # 상품 정보 텍스트
    brand = ""
    for coord, val in all_cells.items():
        if "브랜드" in val or "제품명" in val or "상품명" in val:
            brand = val

    products = []
    for coord, val in all_cells.items():
        row_num = int(re.findall(r'\d+', coord)[0])
        row_vals = {c: v for c, v in all_cells.items() if int(re.findall(r'\d+', c)[0]) == row_num}
        row_text = " | ".join(row_vals.values())
        if len(row_text) > 5 and row_text not in products:
            products.append(row_text)

    full_text = f"브랜드: {brand}\n\n" + "\n".join(products[:80])
    if len(full_text) > 3000:
        full_text = full_text[:3000]

    wb.close()
    return {"text": full_text, "urls": urls, "brand": brand}


def fetch_product_info(url):
    """상품 URL에서 정보를 가져온다"""
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "ko-KR,ko;q=0.9,en;q=0.8",
            "Accept-Encoding": "gzip, deflate",
            "Connection": "keep-alive",
        }
        resp = requests.get(url, headers=headers, timeout=20, verify=False)
        resp.raise_for_status()
        resp.encoding = resp.apparent_encoding or 'utf-8'

        soup = BeautifulSoup(resp.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header", "iframe", "noscript"]):
            tag.decompose()

        title = soup.title.string.strip() if soup.title and soup.title.string else ""

        meta_desc = ""
        meta = soup.find("meta", attrs={"name": "description"})
        if meta and meta.get("content"):
            meta_desc = meta["content"].strip()

        og_data = {}
        for og in soup.find_all("meta", attrs={"property": re.compile(r"^og:")}):
            key = og.get("property", "").replace("og:", "")
            val = og.get("content", "")
            if val:
                og_data[key] = val

        price_texts = []
        for el in soup.find_all(string=re.compile(r'[\d,]+원')):
            text = el.strip()
            if text and len(text) < 100:
                price_texts.append(text)

        main = soup.find("main") or soup.find("article") or soup.find("div", {"id": re.compile(r"content|product|detail", re.I)}) or soup.find("div", {"class": re.compile(r"content|product|detail", re.I)})
        body_text = main.get_text(separator="\n", strip=True) if main else (soup.body.get_text(separator="\n", strip=True) if soup.body else "")

        lines = [l.strip() for l in body_text.split("\n") if l.strip() and len(l.strip()) > 1]
        body_text = "\n".join(lines[:150])

        result = f"[상품 페이지 정보]\n제목: {title}\n"
        if meta_desc:
            result += f"설명: {meta_desc}\n"
        for k, v in og_data.items():
            result += f"{k}: {v}\n"
        if price_texts:
            result += f"가격 정보: {', '.join(price_texts[:5])}\n"
        result += f"\n[상세 내용]\n{body_text[:3000]}"
        return result

    except Exception as e:
        return f"(URL에서 정보를 가져오지 못했습니다: {e})"


def create_storyboard_ppt(sb_data, topic, tone):
    """스토리보드 PPT — 메리쏘드 PPT 구조 정밀 복제"""
    from pptx.oxml.ns import qn
    from pptx.util import Emu
    from copy import deepcopy

    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)

    C_DARK = RGBColor(0x2D, 0x2D, 0x3F)
    C_ACCENT = RGBColor(0xE8, 0x5D, 0x26)
    C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    C_BLACK = RGBColor(0x33, 0x33, 0x33)
    C_GRAY = RGBColor(0x66, 0x66, 0x66)

    def _to_s(v):
        if isinstance(v, dict):
            return "\n".join([f"{k}: {vl}" for k, vl in v.items()])
        if isinstance(v, list):
            return "\n".join([f"- {x}" if isinstance(x, str) else str(x) for x in v])
        return str(v) if v else ""

    def _set_cell(cell, text, font_size=9, bold=False, color=C_BLACK, align=PP_ALIGN.LEFT):
        """셀에 텍스트+스타일 적용"""
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        lines = str(text).split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.alignment = align
            p.space_before = Pt(1)
            p.space_after = Pt(1)
        cell.vertical_anchor = MSO_ANCHOR.TOP

    def _set_cell_bg(cell, hex_color):
        """셀 배경색 (가장 안전한 방식)"""
        tcPr = cell._tc.get_or_add_tcPr()
        for sf in tcPr.findall(qn('a:solidFill')):
            tcPr.remove(sf)
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': hex_color})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

    def _add_header_textbox(slide, text):
        """메리쏘드와 동일: 상단 헤더 텍스트박스 (left=61625, top=~0)"""
        # 헤더 배경 사각형
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(9144000), Emu(580500))
        rect.fill.solid()
        rect.fill.fore_color.rgb = C_DARK
        rect.line.color.rgb = C_DARK
        rect.line.width = Emu(0)
        # 헤더 텍스트
        tf = rect.text_frame
        tf.margin_left = Emu(200000)
        tf.margin_top = Emu(120000)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_WHITE

    def _add_content_table(slide, rows_data, col_widths):
        """메리쏘드와 동일 위치/크기의 4열 테이블 생성"""
        n_rows = len(rows_data)
        n_cols = len(rows_data[0])
        total_w = sum(col_widths)

        # 메리쏘드 동일 위치: left=237900, top=665375
        tbl_shape = slide.shapes.add_table(n_rows, n_cols, Emu(237900), Emu(665375), Emu(total_w), Emu(4300000))
        table = tbl_shape.table

        # 테이블 스타일 초기화
        tbl_el = table._tbl
        tblPr = tbl_el.find(qn('a:tblPr'))
        if tblPr is not None:
            for attr in ['firstRow', 'bandRow', 'lastRow', 'firstCol', 'lastCol']:
                tblPr.attrib.pop(attr, None)

        # 열 너비
        for i, w in enumerate(col_widths):
            table.columns[i].width = Emu(w)

        # 헤더 행 높이
        table.rows[0].height = Emu(380000)

        # 본문 행 높이 (남은 공간 균등 배분)
        body_rows = n_rows - 1
        if body_rows > 0:
            remaining = 4300000 - 380000
            row_h = remaining // body_rows
            for r in range(1, n_rows):
                table.rows[r].height = Emu(row_h)

        # 셀 채우기
        for r_idx, row in enumerate(rows_data):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                if r_idx == 0:
                    # 헤더
                    _set_cell(cell, val, font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
                    _set_cell_bg(cell, '2D2D3F')
                elif c_idx == 0:
                    # 구분 열
                    _set_cell(cell, val, font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
                    _set_cell_bg(cell, '2D2D3F')
                elif c_idx == 1:
                    # 시간/시연 열
                    _set_cell(cell, val, font_size=9, color=C_BLACK, align=PP_ALIGN.CENTER)
                    _set_cell_bg(cell, 'FFFFFF')
                elif c_idx == 2:
                    # 상세 열 (내용 길이에 따라 글자 크기 조절)
                    sz = 9 if len(str(val)) < 300 else 8 if len(str(val)) < 600 else 7
                    _set_cell(cell, val, font_size=sz, color=C_BLACK)
                    _set_cell_bg(cell, 'FFFFFF')
                else:
                    # 기타 열
                    _set_cell(cell, val, font_size=8, color=C_GRAY)
                    _set_cell_bg(cell, 'FFFFFF')

        return tbl_shape

    # 열 너비 (메리쏘드 기준)
    COL_W = [643600, 629075, 5228475, 2167050]  # 구분/시간/상세/기타

    brand = sb_data.get("title", topic)
    platform = sb_data.get("platform", "")
    hosts_name = sb_data.get("hosts", "")
    dur = sb_data.get("total_duration", "")
    scenes = sb_data.get("scenes", [])

    # ════════════════════════════════════════
    # 슬라이드 1: 표지 (메리쏘드 동일)
    # ════════════════════════════════════════
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    # 브랜드명 텍스트 (메리쏘드: left=1506900, top=879125)
    tb = s1.shapes.add_textbox(Emu(1506900), Emu(879125), Emu(6130200), Emu(1327500))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = brand
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_BLACK
    p.alignment = PP_ALIGN.LEFT
    p2 = tf.add_paragraph()
    p2.text = f"{platform} 스토리보드"
    p2.font.size = Pt(16)
    p2.font.color.rgb = C_GRAY
    # 정보 테이블 (메리쏘드: left=1284825, top=2271875)
    info_shape = s1.shapes.add_table(4, 2, Emu(1284825), Emu(2271875), Emu(6574350), Emu(1524000))
    it = info_shape.table
    tblPr = it._tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        for attr in ['firstRow', 'bandRow', 'lastRow', 'firstCol', 'lastCol']:
            tblPr.attrib.pop(attr, None)
    it.columns[0].width = Emu(975400)
    it.columns[1].width = Emu(5598950)
    for ri, (lb, vl) in enumerate([
        ("일시", sb_data.get("broadcast_datetime", datetime.now().strftime("%Y년 %m월 %d일"))),
        ("장소", sb_data.get("broadcast_location", "(미정)")),
        ("진행", hosts_name),
        ("콜타임", "라이브 1시간 전"),
    ]):
        _set_cell(it.cell(ri, 0), lb, font_size=9, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
        _set_cell_bg(it.cell(ri, 0), 'E8E8E8')
        _set_cell(it.cell(ri, 1), vl, font_size=9, color=C_BLACK)
        _set_cell_bg(it.cell(ri, 1), 'FFFFFF')

    # ════════════════════════════════════════
    # 슬라이드 2: 전체 플로우
    # ════════════════════════════════════════
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_textbox(s2, "전체 플로우 - (실제 상황에 따라 상이할 수 있음)")
    n = len(scenes)
    if n > 0:
        cw = min(1500000, 8400000 // n)
        tw = cw * n
        sx = (9144000 - tw) // 2
        ft_shape = s2.shapes.add_table(2, n, Emu(sx), Emu(1500000), Emu(tw), Emu(2500000))
        ft = ft_shape.table
        tblPr = ft._tbl.find(qn('a:tblPr'))
        if tblPr is not None:
            for attr in ['firstRow', 'bandRow', 'lastRow', 'firstCol', 'lastCol']:
                tblPr.attrib.pop(attr, None)
        for c in range(n):
            ft.columns[c].width = Emu(cw)
        ft.rows[0].height = Emu(600000)
        ft.rows[1].height = Emu(1900000)
        for c, sc in enumerate(scenes):
            _set_cell(ft.cell(0, c), str(c+1), font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            _set_cell_bg(ft.cell(0, c), 'E85D26')
            sec = sc.get("section", "")
            _set_cell(ft.cell(1, c), sec, font_size=8, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
            _set_cell_bg(ft.cell(1, c), 'E8EEF7' if c % 2 == 0 else 'FFFFFF')

    concept = sb_data.get("live_concept", "")
    if concept:
        tc = s2.shapes.add_textbox(Emu(300000), Emu(4200000), Emu(8500000), Emu(500000))
        pc = tc.text_frame.paragraphs[0]
        pc.text = f"컨셉: {concept}"
        pc.font.size = Pt(10)
        pc.font.color.rgb = C_GRAY
        pc.font.bold = True

    # ════════════════════════════════════════
    # 슬라이드 3~: 구간별 상세 (메리쏘드 4열 테이블)
    # ════════════════════════════════════════
    for sc in scenes:
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        section = sc.get("section", "")
        _add_header_textbox(sl, f"방송 정보 - {section}")

        script = _to_s(sc.get("host_script", ""))
        screen = _to_s(sc.get("screen_display", ""))
        prod = _to_s(sc.get("product_info", ""))
        direction = _to_s(sc.get("direction_note", ""))
        viewer = _to_s(sc.get("viewer_action", ""))

        # 상세 열 조합
        detail_parts = []
        if script:
            detail_parts.append(script)
        if prod and prod != script:
            detail_parts.append(f"\n[제품/혜택]\n{prod}")
        if screen:
            detail_parts.append(f"\n[화면 표시]\n{screen}")
        detail = "\n".join(detail_parts)

        # 기타 열
        etc_parts = []
        if direction:
            etc_parts.append(direction)
        if viewer:
            etc_parts.append(f"시청자 유도:\n{viewer}")
        etc = "\n".join(etc_parts)

        # 구분 열 줄바꿈
        sec_display = "\n".join(section.split(" ")) if len(section) > 6 else section

        # 헤더 결정: 제품/사은품 슬라이드면 '시연', 나머지는 '시간'
        is_product = any(kw in section for kw in ["제품", "소구", "사은품"])
        header_2 = "시연" if is_product else "시간"

        # 시연 열: 제품이면 시연 방법, 아니면 시간
        col2_val = ""
        if is_product:
            # direction_note에서 시연 관련 내용 추출
            if "손등" in direction or "얼굴" in direction or "입술" in direction:
                demo_parts = []
                for kw in ["손등", "얼굴", "입술", "팔", "손목"]:
                    if kw in direction:
                        demo_parts.append(kw)
                col2_val = "\n/\n".join(demo_parts) if demo_parts else "손등\n/\n얼굴"
            else:
                col2_val = "손등\n/\n얼굴"
        else:
            col2_val = sc.get("duration", "")

        rows = [
            ["구분", header_2, "상세", "기타"],
            [sec_display, col2_val, detail, etc],
        ]

        _add_content_table(sl, rows, COL_W)

    # ════════════════════════════════════════
    # 마지막: 이벤트 & 클로징 정보
    # ════════════════════════════════════════
    sl_e = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_textbox(sl_e, "라이브 컨셉 & 핵심 혜택")

    benefits = sb_data.get("key_benefits", [])
    bt = "\n".join([f"{i+1}. {b}" for i, b in enumerate(benefits)])

    rows = [
        ["항목", "", "내용", ""],
        ["라이브\n컨셉", "", sb_data.get("live_concept", ""), ""],
        ["핵심\n혜택", "", bt, ""],
        ["진행자", "", hosts_name, ""],
        ["톤앤매너", "", tone, ""],
    ]
    _add_content_table(sl_e, rows, COL_W)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ─────────────────────────────────────────
# 페이지 설정
# ─────────────────────────────────────────
st.set_page_config(page_title="STUDIO", page_icon="🎬", layout="wide")

# ─────────────────────────────────────────
# CSS
# ─────────────────────────────────────────
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Noto+Sans+KR:wght@300;400;500;700;900&display=swap');

    .stApp {
        background: #F5F5F5;
        font-family: 'Noto Sans KR', sans-serif;
    }
    header[data-testid="stHeader"] { background: transparent !important; }
    h1, h2, h3, h4, p, span, label, div {
        font-family: 'Noto Sans KR', sans-serif !important;
    }

    /* 헤더 */
    .header-bar {
        background: #FFFFFF;
        border-bottom: 1px solid rgba(255,255,255,0.06);
        padding: 1.2rem 2rem;
        margin: -1rem -1rem 2rem -1rem;
        display: flex;
        align-items: center;
        justify-content: space-between;
    }
    .header-title {
        color: #222;
        font-size: 1.2rem;
        font-weight: 800;
        letter-spacing: -0.5px;
    }
    .header-accent {
        color: #8B5CF6;
    }
    .header-nav {
        display: flex;
        gap: 24px;
    }
    .header-nav-item {
        color: #888;
        font-size: 0.8rem;
        font-weight: 500;
        padding: 6px 14px;
        border-radius: 8px;
        transition: all 0.3s;
    }
    .header-nav-active {
        color: #333333;
        background: rgba(138,92,246,0.1);
    }

    /* 섹션 */
    .section-label {
        color: #8B5CF6;
        font-size: 0.7rem;
        font-weight: 700;
        text-transform: uppercase;
        letter-spacing: 2px;
        margin-bottom: 0.8rem;
    }

    /* 입력 */
    .stTextInput input, .stTextArea textarea {
        background: #FFFFFF !important;
        border: 1px solid #D0D0D0 !important;
        border-radius: 10px !important;
        color: #333333 !important;
        transition: border-color 0.3s !important;
    }
    .stTextInput input:focus, .stTextArea textarea:focus {
        border-color: #8B5CF6 !important;
        box-shadow: 0 0 0 1px rgba(138,92,246,0.2) !important;
    }
    [data-testid="stSelectbox"] > div > div {
        background: #FFFFFF !important;
        border: 1px solid #D0D0D0 !important;
        border-radius: 10px !important;
        color: #333333 !important;
    }
    [data-testid="stDateInput"] > div > div {
        background: #FFFFFF !important;
        border: 1px solid #D0D0D0 !important;
        border-radius: 10px !important;
    }
    [data-testid="stDateInput"] input { color: #333333 !important; }
    label, .stTextInput label, .stSelectbox label,
    [data-testid="stWidgetLabel"] p {
        color: #555555 !important;
        font-size: 0.85rem !important;
    }
    .stMarkdown p, .stMarkdown li { color: #444444 !important; }

    /* 버튼 */
    .stButton > button {
        background: #8B5CF6 !important;
        color: white !important;
        border: none !important;
        padding: 0.75rem 2rem !important;
        font-size: 0.95rem !important;
        font-weight: 700 !important;
        border-radius: 10px !important;
        transition: all 0.3s ease !important;
        box-shadow: 0 4px 12px rgba(138,92,246,0.2);
    }
    .stButton > button:hover {
        background: #7C3AED !important;
        transform: translateY(-2px) !important;
        box-shadow: 0 8px 20px rgba(138,92,246,0.3) !important;
    }
    .stDownloadButton > button {
        background: #06B6D4 !important;
        color: white !important;
        border: none !important;
        padding: 0.75rem 2rem !important;
        font-size: 0.95rem !important;
        font-weight: 700 !important;
        border-radius: 10px !important;
        box-shadow: 0 4px 12px rgba(6,182,212,0.2);
        transition: all 0.3s ease !important;
    }
    .stDownloadButton > button:hover {
        transform: translateY(-2px) !important;
        box-shadow: 0 8px 20px rgba(6,182,212,0.3) !important;
    }

    /* 업로더 */
    [data-testid="stFileUploader"] {
        background: #FFFFFF;
        border: 1px dashed #C0C0C0;
        border-radius: 10px;
        padding: 1.2rem;
    }
    [data-testid="stFileUploader"]:hover {
        border-color: rgba(138,92,246,0.3);
    }

    /* Expander */
    .streamlit-expanderHeader {
        background: #FFFFFF !important;
        border-radius: 10px !important;
        color: #444444 !important;
        font-weight: 500 !important;
    }
    .streamlit-expanderContent {
        background: rgba(255,255,255,0.015) !important;
        border-radius: 0 0 10px 10px !important;
    }

    /* 알림 */
    .stAlert { border-radius: 10px !important; }

    /* 구분선 */
    hr { border-color: rgba(255,255,255,0.05) !important; }

    /* 스크롤바 */
    ::-webkit-scrollbar { width: 4px; }
    ::-webkit-scrollbar-track { background: transparent; }
    ::-webkit-scrollbar-thumb { background: #333; border-radius: 2px; }

    .footer-3d {
        text-align: center;
        padding: 2rem 0;
        color: #333;
        font-size: 0.7rem;
        letter-spacing: 2px;
    }
</style>
""", unsafe_allow_html=True)

# ─── 헤더 ───
st.markdown("""
<div class="header-bar">
    <div class="header-title"><span class="header-accent">LIVE</span> STUDIO</div>
    <div class="header-nav">
        <span class="header-nav-item header-nav-active">스토리보드</span>
        <span class="header-nav-item">배너 생성</span>
        <span class="header-nav-item">분석</span>
    </div>
</div>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────
# 입력 폼
# ─────────────────────────────────────────

# 기본 정보
st.markdown('<div class="section-label">기본 정보</div>', unsafe_allow_html=True)

col1, col2 = st.columns([3, 1])
with col1:
    topic = st.text_input("라이브 주제 / 브랜드명", placeholder="예: 닥터지 신제품 런칭")
with col2:
    style = st.selectbox("톤앤매너", ["밝고 친근한", "전문적이고 신뢰감 있는", "활기차고 에너지 넘치는", "고급스럽고 세련된"])

st.markdown("---")

# 방송 정보
st.markdown('<div class="section-label">방송 정보</div>', unsafe_allow_html=True)

col_a, col_b, col_c = st.columns(3)
with col_a:
    broadcast_platform = st.selectbox(
        "플랫폼",
        ["카카오쇼핑라이브", "네이버 쇼핑라이브", "쿠팡 라이브", "SSG 라이브", "11번가 라이브", "기타"],
        key="bc_platform",
    )
    if broadcast_platform == "기타":
        broadcast_platform = st.text_input("플랫폼명", key="bc_plat_custom")
with col_b:
    hosts = st.text_input("쇼호스트", placeholder="예: 이소유 & 조을희", key="bc_hosts")
with col_c:
    broadcast_duration = st.selectbox("방송 시간", ["30분", "60분", "90분", "120분"], index=1, key="bc_duration")

col_d, col_e, col_f = st.columns(3)
with col_d:
    broadcast_date = st.date_input("방송 날짜", key="bc_date")
with col_e:
    broadcast_time = st.text_input("방송 시간대", placeholder="예: 20:30~21:30", key="bc_time")
with col_f:
    broadcast_location = st.text_input("방송 장소", placeholder="예: 강남구 역삼로 216", key="bc_location")

if broadcast_date:
    date_str = broadcast_date.strftime("%Y년 %m월 %d일")
    datetime_str = f"{date_str} {broadcast_time}" if broadcast_time else date_str
else:
    datetime_str = ""

st.markdown("---")

# 제품 소개서
st.markdown('<div class="section-label">제품 소개서</div>', unsafe_allow_html=True)

col_g, col_h = st.columns([1, 1])
with col_g:
    uploaded_excel = st.file_uploader(
        "엑셀 (.xlsx) 업로드 — 상품 정보 + URL 자동 추출",
        type=["xlsx"],
    )
with col_h:
    with st.expander("URL 직접 입력"):
        st.text_area(
            "한 줄에 하나씩",
            placeholder="https://smartstore.naver.com/...",
            height=80,
            key="sb_urls_manual",
        )

st.markdown("---")

# 생성 버튼
generate_btn = st.button(
    "스토리보드 생성하기",
    use_container_width=True,
    disabled=not topic,
)


# ─────────────────────────────────────────
# 생성 프로세스
# ─────────────────────────────────────────
if generate_btn and topic:
    st.markdown("---")

    excel_data = None
    urls_to_fetch = []
    product_info = ""

    # 1) 엑셀 파싱
    if uploaded_excel:
        with st.spinner("제품 소개서 분석 중..."):
            excel_data = parse_excel_products(uploaded_excel.getvalue())
            urls_to_fetch = excel_data["urls"]
            st.success(f"제품 소개서 분석 완료! (URL {len(urls_to_fetch)}개 발견)")

    # 2) 직접 입력 URL
    manual_urls = st.session_state.get("sb_urls_manual", "")
    if manual_urls.strip():
        extra = [u.strip() for u in manual_urls.strip().split("\n") if u.strip()]
        for u in extra:
            if u not in urls_to_fetch:
                urls_to_fetch.append(u)

    # 3) URL에서 상품 정보 수집 (실패해도 계속 진행)
    all_fetched = []
    failed_count = 0
    if urls_to_fetch:
        with st.spinner(f"상품 페이지 {len(urls_to_fetch)}개에서 정보 수집 중..."):
            for i, url in enumerate(urls_to_fetch):
                info = fetch_product_info(url)
                if "가져오지 못했습니다" in info:
                    failed_count += 1
                else:
                    all_fetched.append(f"[상품 {i+1}] {url}\n{info}")
        if all_fetched:
            st.success(f"상품 정보 {len(all_fetched)}개 수집 완료!")
        if failed_count > 0:
            st.info(f"URL {failed_count}개는 접근이 차단된 사이트입니다. 소개서 내용만으로 진행합니다.")

    # 4) 상품 정보 조합
    parts = []
    if excel_data and excel_data["text"]:
        parts.append(f"[제품 소개서 내용]\n{excel_data['text']}")
    if all_fetched:
        parts.append("\n\n---\n\n".join(all_fetched))
    product_info = "\n\n===\n\n".join(parts)

    # 5) AI 스토리보드 생성
    with st.spinner("AI가 스토리보드를 구성하고 있습니다..."):
        product_context = ""
        if product_info:
            product_context = f"""

아래는 제품 소개서와 실제 상품 페이지에서 가져온 정보야.
각 상품의 실제 이름, 가격, 특징, 장점, 셀링포인트 등을 대본과 자막에 구체적으로 반영해줘.
상품이 여러 개면 각 상품을 자연스럽게 소개하는 장면을 포함해줘.

{product_info}
"""

        sb_prompt = f"""라이브커머스 방송 스토리보드를 JSON으로 작성해.

[입력 정보]
주제/브랜드: {topic}
톤앤매너: {style}
플랫폼: {broadcast_platform}
쇼호스트: {hosts or '(미정)'}
방송 일시: {datetime_str or '(미정)'}
방송 장소: {broadcast_location or '(미정)'}
방송 시간: {broadcast_duration}
{product_context}

[규칙]
- 반드시 아래 순서로 scene을 만들어. 각 구간이 별도 scene이야:

  1. "오프닝 & 브랜드 소개" (1개 scene)
     - 쇼호스트 인사, 오늘 방송 소개, 브랜드 수상이력/신뢰도
     - host_script에 오프닝 대본 + 브랜드 소개 대본을 함께 작성

  2. "혜택 소개" (1개 scene)
     - 라이브 혜택 전체 (쿠폰, 사은품, 이벤트 등)
     - 모든 혜택을 번호로 상세히 나열

  3. "가격 소개" (1개 scene)
     - 모든 제품의 정상가/할인가/쿠폰적용가/할인율
     - 구성품(본품+리필+퍼프 등)도 포함

  4. "사은품 소구 포인트" (1개 scene)
     - 사은품 제품들의 특징, 소구포인트
     - 사은품별로 #해시태그 + 소구포인트 나열

  5. "[제품명] 제품 소개" (제품마다 각각 1개 scene)
     - 소구포인트 (성분, 효능, 특징을 구체적으로)
     - 색상별 추천 피부톤
     - "이런 분들께 추천" 항목
     - 기능성 (SPF, PA 등)
     - direction_note에 시연 방법 (손등시연/얼굴시연/입술시연 등 제품에 맞게)

  6. "이벤트 당첨자 발표 & 클로징" (1개 scene)
     - 이벤트 당첨자 발표 + 혜택 리마인드 + 마무리 인사

- host_script는 실제 방송에서 바로 읽을 수 있는 대화체, 최소 300자 이상
- 제품 소개 scene의 host_script는 500자 이상으로 상세하게
- direction_note의 시연은 제품 특성에 맞게 (쿠션=손등/얼굴, 틴트=입술, 클렌징=세정 시연 등)
- 반드시 JSON만 출력. 설명 텍스트 금지.

```json
{{
  "title": "방송 제목",
  "platform": "{broadcast_platform}",
  "total_duration": "{broadcast_duration}",
  "hosts": "{hosts or '(미정)'}",
  "broadcast_datetime": "{datetime_str or '(미정)'}",
  "broadcast_location": "{broadcast_location or '(미정)'}",
  "scenes": [
    {{
      "scene_number": 1,
      "section": "오프닝 & 브랜드 소개",
      "duration": "",
      "host_script": "쇼호스트 대본 (대화체, 최소 300자)",
      "screen_display": "화면 표시 내용",
      "product_info": "제품/혜택 정보",
      "direction_note": "연출/시연 지시",
      "viewer_action": "시청자 유도"
    }}
  ],
  "live_concept": "컨셉 한줄",
  "key_benefits": ["혜택1", "혜택2", "혜택3"]
}}
```"""

        sb_result = call_claude(sb_prompt)

        # JSON 파싱 (여러 방법으로 시도)
        sb_data = None
        # 방법 1: ```json ... ``` 블록 찾기
        code_match = re.search(r'```(?:json)?\s*([\s\S]*?)```', sb_result)
        if code_match:
            try:
                sb_data = json.loads(code_match.group(1).strip())
            except json.JSONDecodeError:
                pass

        # 방법 2: 전체에서 { } 찾기
        if sb_data is None:
            try:
                # 가장 바깥쪽 { } 찾기
                start = sb_result.index('{')
                depth = 0
                end = start
                for i in range(start, len(sb_result)):
                    if sb_result[i] == '{':
                        depth += 1
                    elif sb_result[i] == '}':
                        depth -= 1
                        if depth == 0:
                            end = i + 1
                            break
                sb_data = json.loads(sb_result[start:end])
            except (ValueError, json.JSONDecodeError):
                pass

        # 방법 3: 실패 시 AI에게 다시 요청
        if sb_data is None:
            with st.spinner("스토리보드를 다시 생성하고 있습니다..."):
                retry_prompt = f"아래 내용을 올바른 JSON 형식으로만 변환해줘. 다른 텍스트 없이 JSON만 출력해줘:\n\n{sb_result[:3000]}"
                retry_result = call_claude(retry_prompt)
                try:
                    retry_match = re.search(r'\{[\s\S]*\}', retry_result)
                    if retry_match:
                        sb_data = json.loads(retry_match.group())
                except (ValueError, json.JSONDecodeError):
                    pass

        if sb_data is None:
            st.error("스토리보드 생성에 실패했습니다. 다시 시도해주세요.")
            with st.expander("AI 응답 원문 (디버깅용)"):
                st.code(sb_result)
            st.stop()

    # 6) 미리보기
    st.markdown(f"""
    <div class="glass-card">
        <div style="color: #8B5CF6; font-size: 0.7rem; font-weight: 700; text-transform: uppercase; letter-spacing: 2px; margin-bottom: 0.5rem;">LIVE COMMERCE STORYBOARD</div>
        <div style="color: #e0e0f0; font-size: 1.3rem; font-weight: 800;">{sb_data.get('title', topic)}</div>
        <div style="color: #777777; font-size: 0.85rem; margin-top: 4px;">{sb_data.get('platform', broadcast_platform)} | {sb_data.get('total_duration', broadcast_duration)} | {sb_data.get('hosts', hosts)} | {len(sb_data.get('scenes', []))}개 구간</div>
    </div>
    """, unsafe_allow_html=True)

    # 컨셉 & 핵심 혜택
    live_concept = sb_data.get("live_concept", "")
    key_benefits = sb_data.get("key_benefits", [])
    if live_concept or key_benefits:
        with st.expander("라이브 컨셉 & 핵심 혜택", expanded=True):
            if live_concept:
                st.markdown(f"**컨셉:** {live_concept}")
            if key_benefits:
                st.markdown("**핵심 혜택:**")
                for b in key_benefits:
                    st.markdown(f"- {b}")

    def _preview_str(v):
        if isinstance(v, dict):
            return "\n".join([f"- {k}: {val}" for k, val in v.items()])
        if isinstance(v, list):
            return "\n".join([f"- {item}" for item in v])
        return str(v) if v else ""

    scenes = sb_data.get("scenes", [])
    for scene in scenes:
        with st.expander(f"구간 {scene['scene_number']}  —  {scene.get('section', '')}  ({scene.get('duration', '')})"):
            st.markdown(f"**쇼호스트 대본:**")
            st.markdown(_preview_str(scene.get('host_script', '')))
            st.markdown("---")
            col_a, col_b = st.columns(2)
            with col_a:
                st.markdown(f"**화면 표시:** {_preview_str(scene.get('screen_display', ''))}")
                st.markdown(f"**제품/혜택 정보:** {_preview_str(scene.get('product_info', ''))}")
            with col_b:
                st.markdown(f"**연출 지시:** {_preview_str(scene.get('direction_note', ''))}")
                st.markdown(f"**시청자 유도:** {_preview_str(scene.get('viewer_action', ''))}")

    # 7) PPT 다운로드
    with st.spinner("PPT 파일 생성 중..."):
        ppt_bytes = create_storyboard_ppt(sb_data, topic, style)

    st.success("라이브커머스 스토리보드 PPT가 생성되었습니다!")

    filename = f"라이브커머스_스토리보드_{topic[:15]}_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    st.download_button(
        label=f"PPT 다운로드  —  {filename}",
        data=ppt_bytes,
        file_name=filename,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
    )


# ─────────────────────────────────────────
# 푸터
# ─────────────────────────────────────────
st.markdown(
    '<div class="footer-3d">STUDIO v1.0 &nbsp;&middot;&nbsp; Powered by Claude AI &nbsp;&middot;&nbsp; '
    f'{datetime.now().strftime("%Y.%m.%d")}</div>',
    unsafe_allow_html=True,
)
